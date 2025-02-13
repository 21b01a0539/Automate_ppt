import streamlit as st
from openai import OpenAI
import PyPDF2
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
import os
import re
from dotenv import load_dotenv

load_dotenv()

def get_openai_client():
    """
    Initialize OpenAI client with a default API key
    """
    # Using a default API key
    openai_api_key = "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"
    return OpenAI(api_key=openai_api_key)

def extract_pdf_text(uploaded_file):
    """
    Extract text from an uploaded PDF file.
    """
    try:
        pdf_reader = PyPDF2.PdfReader(uploaded_file)
        
        full_text = ""
        for page in pdf_reader.pages:
            full_text += page.extract_text() + "\n\n"
        
        return full_text
    except Exception as e:
        st.error(f"Error extracting PDF text: {e}")
        return ""

def generate_slide_content(client, text):
    """
    Use OpenAI to generate slide content with improved parsing strategy.
    """
    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert presentation designer converting complex research papers into clear, concise slides."},
                {"role": "user", "content": f"""Convert the research paper text into a structured presentation. 
                Follow these strict formatting guidelines:
                - TITLE FORMAT: [Slide Title]
                - BULLETS FORMAT: 
                  • First bullet point
                  • Second bullet point
                  • Third bullet point
                  ...(As many are needed for the slide (atmost 5 bullet points))

                Slide Structure:
                1. Title Slide
                2. Research Background/Introduction
                3. Methodology
                4. Key Findings
                5. Results and Implications
                6. Conclusion
                7. (Optional) Future Research Directions

                Formatting Rules:
                - Use exact formatting as shown above
                - Keep titles short (5-7 words)
                - 3-4 concise bullet points per slide
                - Avoid complex jargon

                Research Text (first 4000 characters):
                {text[:4000]}"""}
            ],
            max_tokens=1500
        )
        
        # Get the raw content
        raw_content = response.choices[0].message.content
        
        # Advanced parsing with robust error handling
        def parse_slides(content):
            # Split by lines and clean
            lines = [line.strip() for line in content.split('\n') if line.strip()]
            
            slides = []
            current_slide = []
            
            for line in lines:
                # Check if line is a potential title (starts with capital letter, not a bullet)
                if re.match(r'^[A-Z][a-z\s]+', line) and not line.startswith('•'):
                    # If we have a previous slide, add it
                    if current_slide:
                        slides.append('\n'.join(current_slide))
                        current_slide = []
                    
                    # Start new slide with this title
                    current_slide.append(line)
                elif line.startswith('•'):
                    # Bullet point
                    current_slide.append(line)
            
            # Add last slide
            if current_slide:
                slides.append('\n'.join(current_slide))
            
            return slides
        
        # Parse slides with error handling
        parsed_slides = parse_slides(raw_content)
        
        # Filter out empty or invalid slides
        valid_slides = [slide for slide in parsed_slides if len(slide.split('\n')) > 1]
        
        return valid_slides
    
    except Exception as e:
        st.error(f"Error generating slide content: {e}")
        return []

def create_powerpoint(slide_contents):
    """
    Create a PowerPoint presentation with robust parsing.
    """
    prs = Presentation()
    
    # Color palette
    COLOR_PALETTE = {
        'background': RGBColor(240, 248, 255),
        'title': RGBColor(0, 51, 102),
        'text': RGBColor(51, 51, 51),
    }

    # Slide layouts
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    
    # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Research Insights"
    subtitle.text = "Comprehensive Summary"
    
    # Content slides
    for content in slide_contents:
        # Robust parsing of slide content
        lines = content.split('\n')
        
        # First line is the title, rest are bullets
        slide_title = lines[0].strip()
        slide_bullets = [line.strip().lstrip('•').strip() for line in lines[1:] if line.strip()]
        
        # Add slide
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Customize slide title
        title = slide.shapes.title
        title.text = slide_title
        
        # Add bullet points
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        for bullet in slide_bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
    
    # Save to in-memory file
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    return pptx_buffer

def main():
    st.title("Research Paper to PowerPoint Converter")
    
    # Color scheme selection
    st.sidebar.header("Presentation Customization")
    color_scheme = st.sidebar.selectbox(
        "Choose Color Scheme", 
        ["Professional Blue", "Corporate Gray", "Modern Minimal"]
    )
    
    # Get OpenAI Client
    client = get_openai_client()
    
    if not client:
        return
    
    # File uploader
    uploaded_file = st.file_uploader("Upload a Research Paper (PDF)", type=['pdf'])
    
    if uploaded_file is not None:
        # Extract text
        with st.spinner('Extracting text from PDF...'):
            extracted_text = extract_pdf_text(uploaded_file)
        
        if extracted_text:
            # Generate slide content
            with st.spinner('Generating slide content with AI...'):
                slide_contents = generate_slide_content(client, extracted_text)
            
            if slide_contents:
                # Preview slide contents
                st.subheader("Generated Slide Contents")
                for i, slide in enumerate(slide_contents, 1):
                    st.text(f"Slide {i}: {slide}")
                
                # Create PowerPoint
                with st.spinner('Creating PowerPoint presentation...'):
                    pptx_file = create_powerpoint(slide_contents)
                
                # Download button
                st.download_button(
                    label="Download PowerPoint Presentation",
                    data=pptx_file,
                    file_name="research_paper_summary.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            else:
                st.warning("Could not generate slide content.")
        else:
            st.warning("Could not extract text from the PDF.")

if __name__ == "__main__":
    main()