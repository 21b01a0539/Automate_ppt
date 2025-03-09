from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.dml.fill import FillFormat
import streamlit as st
import io
from components import get_relevant_image

def create_ppt(slides_content, heading_rgb, heading_size, bg_rgb, content_rgb, content_size, heading_font, content_font, api_key=None):
    """Create PowerPoint presentation with slides and optional images"""
    prs = Presentation()
    
    first_slide = True
    for title, content_list in slides_content.items():
        try:
            print(f"\nProcessing slide: {title}")  # Debug print
            
            if first_slide:
                # First slide with layout 0 (Title Slide)
                
                slide_layout = prs.slide_layouts[0]  # Title Slide Layout
                slide = prs.slides.add_slide(slide_layout)

                title_shape = slide.shapes.title
                subtitle = slide.placeholders[1]

                # Assign title
                title_shape.text = content_list[0]  # First line as title
    
                # Assign subtitle (Joining remaining lines)
                subtitle.text = "\n".join(content_list[1:])  # Joining rest of the content
                

  
                for para in subtitle.text_frame.paragraphs:
                    para.font.size = Pt(18)  # Adjust the size as needed

                first_slide = False  # Mark first slide as processed
            else:
                # Other slides with layout 1 (Title & Content)
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title and content
                title_shape = slide.shapes.title
                title_shape.text = title
                content = slide.placeholders[1]
                content.text = "\n".join([f"{point}" for point in content_list])
                
                # Try to add a relevant image
                try:
                    print(f"Attempting to fetch image for: {title}")
                    image_data = get_relevant_image(title)
                    if image_data:
                        print("Successfully got image data")
                        # Add image to right side of slide
                        left = prs.slide_width - 4000000  # Align to the right
                        top = prs.slide_height - 3100000  # Align to the bottom

                        width = 3500000  # Adjust size
                        height = 2500000  # Adjust size
                        slide.shapes.add_picture(image_data, left, top, width=width, height=height)
                        print(f"Image added to slide: {title}")
                    else:
                        print(f"No image data received for: {title}")
                except Exception as img_error:
                    print(f"Error adding image to slide: {str(img_error)}")
                
                # Apply styling for content slides
                for p in content.text_frame.paragraphs:
                    p.font.size = Pt(content_size)
                    p.font.color.rgb = RGBColor(*content_rgb)
                    p.font.name = content_font
            
            # Apply styling for all slides
            title_shape.text_frame.paragraphs[0].font.size = Pt(heading_size)
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*heading_rgb)
            title_shape.text_frame.paragraphs[0].font.name = heading_font
            
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_rgb)
            
        except Exception as slide_error:
            print(f"Error processing slide {title}: {str(slide_error)}")

    # Save presentation
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer


def main():
    st.title("Automated PowerPoint Presentation Generator")

    # Sidebar for customization options
    st.sidebar.header("Customize Your Presentation")
    heading_color = st.sidebar.color_picker("Choose Heading Color", "#000000")
    heading_size = st.sidebar.slider("Heading Font Size", 20, 60, 30)
    heading_font = st.sidebar.selectbox("Select Heading Font", 
                                        ["Arial", "Helvetica", "Times New Roman", "Courier New", "Georgia", "Verdana", 
                                         "Trebuchet MS", "Calibri", "Cambria", "Garamond", "Comic Sans MS", 
                                         "Palatino Linotype", "Tahoma", "Lucida Sans Unicode", "Impact", 
                                         "Franklin Gothic Medium", "Segoe UI", "Optima", "Baskerville", "Bookman","Algerian"])
    content_color = st.sidebar.color_picker("Choose Content Color", "#000000")
    content_size = st.sidebar.slider("Content Font Size", 10, 40, 20)
    content_font = st.sidebar.selectbox("Select Content Font", 
                                        ["Arial", "Helvetica", "Times New Roman", "Courier New", "Georgia", "Verdana", 
                                         "Trebuchet MS", "Calibri", "Cambria", "Garamond", "Comic Sans MS", 
                                         "Palatino Linotype", "Tahoma", "Lucida Sans Unicode", "Impact", 
                                         "Franklin Gothic Medium", "Segoe UI", "Optima", "Baskerville", "Bookman","Algerian"])
    bg_color = st.sidebar.color_picker("Choose Background Color", "#FFFFFF")
    unsplash_access_key = st.sidebar.text_input("Unsplash Access Key")

    # Convert hex to RGB
    heading_rgb = tuple(int(heading_color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
    content_rgb = tuple(int(content_color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
    bg_rgb = tuple(int(bg_color.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))

    # Define the content
    content = {
        "Title": [
            "Automated Presentation Slide Generation from Research Papers using NLP and Deep Learning",
            "Atul Shreewastav, Bidhan Acharya, Nischal Paudel, Yugratna Humagain",
            "Department of Electronics and Computer Engineering, IOE, Thapathali Campus",
            "May 13, 2024"
        ],
        "Abstract": [
            "Novel approach for automating generation of presentation slides from academic research papers",
            "Leveraging NLP techniques, specifically a fine-tuned T5 transformer model",
            "Custom dataset of research articles used for training the model",
            "Enables automatic conversion into concise and informative slides",
            "Potential impact on summarizing and communicating technical information within research community"
        ],
        "Introduction": [
            "Communication of complex information in professional and educational settings relies on clear presentations",
            "Manual creation of presentations time-consuming and resource-intensive",
            "NLP advancements offer opportunities for automating content extraction and organization",
            "Proposal of automated slide generation methodology using fine-tuned T5 transformer model",
            "Aim to empower users to focus on content creation and knowledge dissemination"
        ],
        "Methodology": [
            "Challenges associated with manual slide creation addressed using NLP techniques",
            "Fine-tuned T5 transformer model utilized for extracting key information and summarizing content",
            "Model trained on custom dataset of research articles for proficiency in text-to-text transfer tasks",
            "Approach enables automatic conversion of research articles into well-structured presentation slides"
        ],
        "Related Work": [
            "Active research area in automating generation of presentation slides from research papers",
            "Examples of existing approaches like PPSGen, Learning Based Slide Generator, and DOC2PPT",
            "Emphasis on factual accuracy, machine learning techniques, and analyzing paper structure for slide generation",
            "Various methods explored for automatic slide generation from scientific documents"
        ],
        "Results": [
            "Successful implementation of automated slide generation methodology",
            "Extraction of key information and creation of concise slides demonstrated",
            "Comparison of manually created slides with those generated automatically",
            "Improved efficiency and clarity in presentation creation process observed"
        ],
        "Conclusions": [
            "Novel approach using NLP and deep learning techniques for automating slide generation from research papers",
            "Potential to revolutionize how technical information is summarized and communicated",
            "Future directions include enhancing the model's capabilities and expanding to other domains",
            "Contributions to efficient and effective presentation creation process for professional and educational contexts"
        ]
    }

    # Generate button
    if st.button("Generate Presentation"):
        pptx_file = create_ppt(content, heading_rgb, heading_size, bg_rgb, content_rgb, content_size, heading_font, content_font, unsplash_access_key)
        st.download_button(label="Download Presentation", data=pptx_file, file_name="generated_presentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        st.success("Presentation generated successfully!")

if __name__ == "_main_":
    main()
